import numpy as np
import cv2
from PIL import ImageTk
import imutils
from imutils.object_detection import non_max_suppression
from imutils import paths
import pptx
import pptx.util
from pptx.util import Inches, Pt
import tkinter
from tkinter import *
from tkinter import filedialog
from tkinter.filedialog import askdirectory
import sys
import os
import time
import math
import random
import re

imgW = 400 #400
noC = 10 #10
minPts = 6 #6
matchThres = 85
estabThres = 3 #3
decayThres = 4
filterMse = 1400.0
pptName = 'CourseSlides'
argList = ['r','','']
argList2 = ['n','']
txtInS = False
vloaded1=False
vloaded2=False

class cellState:

	def __init__(self,pts,birthF):
		self.pts = pts
		self.birthF = birthF
		self.timeOS = 0
		self.timeD = 0

class Lec2Doc:

	def __init__(self,sampling=30,fps=30):
		global imgW,noC,minPts,matchThres,estabThres,decayThres,\
		argList,argList2,txtInS,filterMse,pptName
		self.argList = argList
		self.argList2 = argList2
		self.txtInS = txtInS
		self.sampling = sampling
		self.filterMse = filterMse
		self.pptName = pptName
		self.fps = 30
		self.imgW = imgW
		self.noC = noC
		minPts = float(minPts)
		self.minPts =  (((imgW/noC)**2)/100) * minPts
		print("MIN PTS")
		print(self.minPts)
		self.matchThres = matchThres
		self.estabThres = estabThres #3
		self.decayThres = decayThres
		self.pastJ = 5
		self.endJ = 4
		self.useExtremas = False

		self.resetData()
		self.state2Col = {0:(155,0,0),1:(255,0,0),2:(255,155,0),3:(0,155,0)\
		,4:(0,255,0),5:(255,255,50),6:(255,255,200)}
		self.state2Col2 = {0:(0,255,255),1:(0,150,255),2:(0,50,200)\
		,3:(0,50,155),4:(0,50,55),5:(0,0,50),6:(0,0,0)}

	def resetData(self,full=True):
		self.curFrames = list()
		self.toSaveFrames = list()
		self.toSaveTimes = list()
		self.captions = list()
		self.cellGrid = list()
		self.tranSeq = list()
		self.prevSedge = 0
		if full:
			self.tranSeqFolder = list()


	def getNewCell(self,edges,r,c,t):
		eroi = self.rc2ROI(edges,r,c)
		croi = self.rc2ROI(self.allContImage,r,c)
		pts = list(np.where(eroi==255) )

		noPts = len(pts[0])
		if noPts < self.minPts:
			return False
		else:
			pts = [  [ pts[0][n], pts[1][n] ]  for n in range( noPts )     ]
			if self.ptsInROI(pts,croi,False)==False:
				cell = cellState(pts,t)
				return cell
			else:
				return False

	def deleteCells(self,toKill):

		tempCellGrid = [[ [ ] for x in range(self.noC)] for y in range(self.noR)]

		for r in range(self.noR):
			for c in range(self.noC):
				curCL = self.cellGrid[r][c]
				for cn,cstate in enumerate( curCL  ):
					if cn not in toKill[r][c]:
						tempCellGrid[r][c].append( cstate  )

		self.cellGrid = tempCellGrid


	def cell4mROI(self,roi,t):
		pts = list(np.where(roi==255) )

		noPts = len(pts[0])
		if noPts < self.minPts:
			return False
		else:
			pts2 = [  [ pts[0][n], pts[1][n] ]  for n in range( noPts )     ]
			cell = cellState(pts2,t)
			return cell

	def rc2ROI(self,img,r,c):
		bound = self.boundsGrid[r][c]
		roi = img[  bound[1]  : bound[3]   ,  bound[0]  : bound[2]    ]
		return roi

	def ptsInROI(self,pts,roi,paint=True):
		found = False
		noPts = float(len(pts))
		matchd = 0.0

		for pt in pts:
			r = pt[0]
			c = pt[1]
			if roi[r][c]==255:
				matchd += 1
			if paint==True:
				self.allContImage[r][c] = 255


		prcntMtch = (matchd * 100.0)/noPts

		if prcntMtch >= self.matchThres:
			found = True

		return found

	def getEdgeImage(self,img,mn,mx):
		gimg = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
		edges = cv2.Canny(gimg, mn,mx)
		return edges


	def getMSE(self,imageA,imageB):
		mse = np.sum((imageA.astype("float") - imageB.astype("float")) ** 2)
		mse /= float(imageA.shape[0] * imageA.shape[1])
		return mse

	def extractFrames(self,path,limit):
		cap = cv2.VideoCapture(path)
		noF = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
		#print("Total No of Frames:")
		#print(noF)
		n=0
		loadFrame=1

		while(cap.isOpened()):
			ret, frame = cap.read()
			if ret==True:
				n += 1
				if n % self.sampling==0:
					self.curFrames.append(frame)
					print("Added Frame No: "+str(loadFrame))

					imgu = np.zeros([int(self.imgW*0.5),int(self.imgW),3],dtype=np.uint8)
					cv2.putText(imgu, "Loading Video Frames : "+str(  int( loadFrame*100/(noF/self.sampling) ) )+" %", \
						(int(self.imgW*0.25), int(self.imgW*0.25)),cv2.FONT_HERSHEY_SIMPLEX, 0.4, (255, 255, 255), lineType=cv2.LINE_AA)
					cv2.imshow("Visualization",imgu)
					cv2.waitKey(1)


					loadFrame += 1

			if (limit != -1 and n > limit) or (limit == -1 and n >= noF):
				break

		cap.release()
		cv2.destroyAllWindows()
		loadFrame = -1

	def runProcess(self,opt):

		self.resetData()

		if opt==1:

			filaname = self.argList[1][5:]
			slsh = filaname.find("/")
			foldiname = filaname[:slsh]

			self.lecture2doc(self.argList[1],self.argList[2])
			self.loadTranscript(foldiname,argList[2])
			self.mergeFramesTranscripts()
			self.generatePPT('Lesson : '+foldiname,'Topic : '+str(argList[2]),self.pptName+'.pptx', self.txtInS   )
		else:

			foldiname = self.argList2[1][5:]
			fullFPath = self.argList2[1]

			directory_list = []
			for root, dirs, files in os.walk(fullFPath, topdown=False):
				for name in dirs:
					directory_list.append(os.path.join(root, name))

			vidFolder = [s for s in directory_list if 'Subtitles' not in s ][0]
			filenames=[]

			for root, dirs, files in os.walk(vidFolder, topdown=False):
				filenames = files

			filenames = [  f for f in filenames if '.' in f[1:]  ]

			atoi = lambda txt: int(txt) if txt.isdigit() else txt
			natKey = lambda txt:[ atoi(c) for c in re.split('(\d+)', txt) ]
			filenames.sort(key=natKey)

			for vidiname in filenames:
				self.resetData(False)
				self.lecture2doc( vidFolder ,  vidiname  )
				self.loadTranscript(foldiname,  vidiname   )
				self.mergeFramesTranscripts()
				self.tranSeqFolder = self.tranSeqFolder + self.tranSeq.copy()

			self.tranSeq = self.tranSeqFolder
			self.generatePPT('Lesson : '+foldiname,' ',self.pptName+'.pptx', self.txtInS)

	def storeSlideImages(self):
		for n,frame in enumerate(self.toSaveFrames):
			cv2.imwrite( "tempImages/img_"+str(n)+".png" , frame  )
			#print("Saved Img: "+str(n))

	def generatePPT(self,titleTxt,sbtText,pptPath,useTxt=False):

		prs = pptx.Presentation()
		#prs.slide_height = 5143500
		slide = prs.slides.add_slide(prs.slide_layouts[0])

		title = slide.shapes.title
		title.text = titleTxt
		subtitle = slide.placeholders[1]
		subtitle.text = sbtText

		if useTxt==False:

			for n,seq in enumerate(self.tranSeq):
				img = seq[1]

				if len(img) != 0:

					slide = prs.slides.add_slide(prs.slide_layouts[6])

					pic_left  = int(prs.slide_width * 0.15)
					pic_top   = int(prs.slide_height * 0.1)
					pic_width = int(prs.slide_width * 0.7)
					pic_height = int(pic_width * img.shape[0] / img.shape[1])

					fName = "tempImages/img_"+str(n)+".png"
					cv2.imwrite( fName , img  )

					pic = slide.shapes.add_picture( fName , pic_left, pic_top, pic_width, pic_height)

					os.remove(fName)

			prs.save(pptPath)
		else:

			for n,seq in enumerate(self.tranSeq):
				txt = seq[0]
				#print(txt)
				ci = 0
				nS = 1
				nC = 1
				cpl = 99
				stl = 16
				lastStart = 0

				while True:
					sch = txt[ci]
					if sch == '\n':
						nC = 0
						nS += 1
					else:
						if nC > cpl:
							j=ci
							while j>0:
								if txt[j]==' ':
									ci=j
									break
								j -= 1
							nC = 0
							nS += 1
					ci += 1
					nC += 1

					if nS > stl:
						toWrite = txt[lastStart:ci]
						slide = prs.slides.add_slide( prs.slide_layouts[6] )
						textBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25),Inches(9.0), Inches(7.0))
						textFrame = textBox.text_frame
						textFrame.word_wrap = True
						textParagraph = textFrame.add_paragraph()
						textParagraph.text = toWrite
						lastStart = ci
						nS = 1
						nC = 1

					if ci>=len(txt):
						toWrite = txt[lastStart:]
						slide = prs.slides.add_slide( prs.slide_layouts[6] )
						textBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25),Inches(9.0), Inches(7.0))
						textFrame = textBox.text_frame
						textFrame.word_wrap = True
						textParagraph = textFrame.add_paragraph()
						textParagraph.text = toWrite
						break

				img = seq[1]

				if len(img) != 0:

					slide = prs.slides.add_slide(prs.slide_layouts[6])

					pic_left  = int(prs.slide_width * 0.15)
					pic_top   = int(prs.slide_height * 0.1)
					pic_width = int(prs.slide_width * 0.7)
					pic_height = int(pic_width * img.shape[0] / img.shape[1])

					fName = "tempImages/img_"+str(n)+".png"
					cv2.imwrite( fName , img  )

					pic = slide.shapes.add_picture( fName , pic_left, pic_top, pic_width, pic_height)
					os.remove(fName)
				

			prs.save(pptPath)

	def lecture2doc(self,lesson,lecture):

		self.extractFrames(lesson+'/'+lecture,-1) #1000

		startComp = time.time()

		for t,frm in enumerate(self.curFrames):
			start = time.time()

			self.frmcp = imutils.resize(frm,width=min(self.imgW,frm.shape[1]))
			frmcp2 = self.frmcp.copy()

			edges = self.getEdgeImage(self.frmcp,30,200)
			cv2.imshow("Edges",edges)
			#cv2.waitKey(0)

			img2, contours, hierarchy = cv2.findContours(  edges ,cv2.RETR_TREE,cv2.CHAIN_APPROX_SIMPLE)
			#print(  len(contours) )
			
			if t>0:
				# State Visualization
				for r in range(self.noR):
					for c in range(self.noC):
						cv2.rectangle(frmcp2,  (self.boundsGrid[r][c][0],self.boundsGrid[r][c][1]), (self.boundsGrid[r][c][2],self.boundsGrid[r][c][3]),(0,255,0),1)
						curCL = self.cellGrid[r][c]
						if len(curCL) > 0:
							for cstate in curCL:
								ctd = cstate.timeOS
								cd = cstate.timeD
								rX = random.randint(self.boundsGrid[r][c][0] + 5  , self.boundsGrid[r][c][2] -5 )
								rY = random.randint(self.boundsGrid[r][c][1] + 5 , self.boundsGrid[r][c][3] - 5)
								cv2.circle(frmcp2, (rX,rY) , 9 , self.state2Col[ctd] , -1)
								cv2.circle(frmcp2, (rX,rY) , 5 , self.state2Col2[cd] , -1)

				gonna = False

				toKill = [[ [ ] for x in range(self.noC)] for y in range(self.noR)]
				#transGs = [[ [ ] for x in range(self.noC)] for y in range(self.noR)]

				self.allContImage = np.zeros(  edges.shape , np.uint8)

				for r in range(self.noR):
					for c in range(self.noC):
						curCL = self.cellGrid[r][c]
						if len(curCL)>0:
							for cn,cstate in enumerate(curCL):
								roi = self.rc2ROI(edges,r,c)
								if self.ptsInROI(cstate.pts, roi):

									if cstate.timeD > 0:
										self.cellGrid[r][c][cn].timeD = 0
									elif cstate.timeD == 0:
										if self.cellGrid[r][c][cn].timeOS < self.estabThres:
											self.cellGrid[r][c][cn].timeOS += 1

								else:

									if cstate.timeD > 0:
										if cstate.timeD < self.decayThres:
											self.cellGrid[r][c][cn].timeD += 1
										else:
											if self.cellGrid[r][c][cn].timeOS < self.estabThres:
												toKill[r][c].append(cn)
											else:
												#transGs[r][c].append(cn)
												gonna = True
												toKill[r][c].append(cn)
									else:
										if self.cellGrid[r][c][cn].timeOS < 2:
											# premature death
											toKill[r][c].append(cn)
										else:
											self.cellGrid[r][c][cn].timeD = 1

				self.deleteCells(toKill)
				estabPres = False

				# check if either no green or green decaying
				

				for r in range(self.noR):
					for c in range(self.noC):
						cell = self.getNewCell(edges,r,c,t)
						curCL = self.cellGrid[r][c]

						NDEstab = False


						for cstate in curCL:

							if cstate.timeOS >= self.estabThres:
								estabPres = True
								#newCheck = True
								#break

								if cstate.timeD == 0:
									NDEstab = True
									break




						if cell != False and NDEstab == False and ((r>0 and c>0 and r < self.noR-1 and c < self.noC-1) or self.useExtremas) :
							self.cellGrid[r][c].append( cell  )

				endCap = estabPres==True and t==len(self.curFrames)-1

				if gonna or endCap:

					#print("Gone at "+str(t))

					capT = t - self.pastJ
					if endCap:
						capT = t - self.endJ

					leena = len(self.toSaveTimes)

					svFr = imutils.resize(self.curFrames[ capT ],width=min(self.imgW,self.curFrames[ capT ].shape[1]))
					svedges = self.getEdgeImage(svFr ,30,200)

					if leena > 0:

						mse = self.getMSE(self.prevSedge, svedges)
						"""
						print("MSE")
						print(mse)
						cv2.imshow("Prev",self.prevSedge)
						cv2.imshow("New",svedges)
						cv2.waitKey(0)
						"""
						if self.toSaveTimes[leena-1] != capT and mse >= self.filterMse :
							self.toSaveFrames.append(  self.curFrames[ capT ]  )
							self.toSaveTimes.append( capT )
							self.prevSedge = svedges.copy()

					else:
						self.toSaveFrames.append(  self.curFrames[ capT ]  )
						self.toSaveTimes.append( capT )
						self.prevSedge = svedges.copy()
					
					toKill = [[ [ ] for x in range(self.noC)] for y in range(self.noR)]

					for r in range(self.noR):
						for c in range(self.noC):
							curCL = self.cellGrid[r][c]

							if len(curCL)>0:
								for cn,cstate in enumerate(curCL):

									if cstate.birthF < capT  or (cstate.timeOS >= self.estabThres and cstate.timeD > 0)  :
										toKill[r][c].append(cn)


					self.deleteCells(toKill)

				#cv2.drawContours( frmcp, contours, -1, (0,255,0), 1)
				#cv2.imshow("Contours Types",self.frmcp)
				cv2.imshow("Visualization",frmcp2)
				#cv2.imshow("Prev",prevEdges)

				cv2.waitKey(1)

			else:

				self.imgH = frmcp2.shape[0]
				self.celLen = float(self.imgW)/float(self.noC)
				self.noR = int(  math.ceil(  self.imgH / self.celLen    )  )
				self.boundsGrid = [[ [0,0,0,0] for x in range(self.noC)] for y in range(self.noR)]

				for r in range(self.noR):
					for c in range(self.noC):
						self.boundsGrid[r][c][0] = int(self.celLen * c)
						self.boundsGrid[r][c][1] = int(self.celLen * r)
						self.boundsGrid[r][c][2] = int(self.celLen * c + self.celLen)
						if self.boundsGrid[r][c][2] >= self.imgW:
							self.boundsGrid[r][c][2] = int(self.imgW-1)
						self.boundsGrid[r][c][3] = int(self.celLen * r + self.celLen)
						if self.boundsGrid[r][c][3] >= self.imgH:
							self.boundsGrid[r][c][3] = int(self.imgH-1)

				self.cellGrid = [[ [ ] for x in range(self.noC)] for y in range(self.noR)]

				for r in range(self.noR):
					for c in range(self.noC):
						roi = self.rc2ROI(edges,r,c)
						cell = self.cell4mROI(roi,t)
						if cell != False and ((r>0 and c>0 and r < self.noR-1 and c < self.noC-1) or self.useExtremas):
							self.cellGrid[r][c].append( cell  )

			prevEdges = edges.copy()
			prevContours = contours.copy()

			end = time.time()
			elapsed = end - start
			#print("Elapsed")
			#print(elapsed)

		endComp = time.time()
		elapsedComp = endComp - startComp
		print("Total Run Time")
		print(elapsedComp)

		cv2.destroyAllWindows()

	def loadTranscript(self,folderPath,videoPath):
		#print("Loading Transcript Data...")

		if self.txtInS: 

			directory_list = []
			for root, dirs, files in os.walk("Data/"+folderPath, topdown=False):
				for name in dirs:
					directory_list.append(os.path.join(root, name))

			subFolder = [s for s in directory_list if 'Subtitles' in s ][0]
			vidFolder = [s for s in directory_list if 'Subtitles' not in s ][0]

			filenames=[]

			for root, dirs, files in os.walk(vidFolder, topdown=False):
				filenames = files

			fileIndex = filenames.index(videoPath)

			for root, dirs, files in os.walk(subFolder, topdown=False):
				filenames = files

			sbtFilename = filenames[fileIndex]
			fullSubtPath = subFolder+'/'+sbtFilename

			file = open(fullSubtPath,'r')
			ftxt = file.read()
			lines=ftxt.splitlines()

			captions = []

			for n in range(len(lines)):

				lineStr = lines[n]

				if lineStr.isdigit():

					ts = lines[n+1]
					starti = ts[0:12]
					endi = ts[17:]

					hS = int( ts[0:2] )
					mS = int( ts[3:5] )
					sS = int( ts[6:8] )
					uS = int( ts[9:12] )

					miliS = hS*3600000 + mS*60000 + sS*1000 + uS

					passage = ''

					ju=n+2

					while ju < len(lines):

						if len( lines[ju] ) <= 2:
							break

						if ju > n+2:
							passage += ' '+lines[ju]
						else:
							passage += lines[ju]
						ju += 1 

					caption = [  miliS, passage  ]
					captions.append(caption)

			self.captions = captions

		else:
			self.captions = []

	def mergeFramesTranscripts(self):
		figImages = self.toSaveFrames
		frameCnts = self.toSaveTimes
		self.tranSeq = []
		lastMS = 0
		capsUsed = 0
		for n in range( len(frameCnts)  ):
			curFig = figImages[n]
			curms = frameCnts[n] * 1000
			#print('Frame '+str(frameCnts[n]))
			#print('MS '+str(curms))

			if curms != lastMS:
				txt = ''
				i = capsUsed
				while i < len(self.captions):
					caption = self.captions[i]

					if caption[0] >= curms:
						capsUsed = i
						break
					else:
						txt += caption[1] + ' '
						i += 1

				self.tranSeq.append( [ txt, curFig  ]   )

			lastMS = curms

		if capsUsed < len(self.captions):
			txt = ''
			i=capsUsed
			while i < len(self.captions):
				caption = self.captions[i]
				txt += caption[1] + ' '
				i += 1
			self.tranSeq.append( [ txt, np.array([])  ]   )

# creating tkinter and canvas objects

root = tkinter.Tk()
root.wm_title("SlideThruLectures")
canvas = Canvas( width = 900, height = 600, bg = 'LightBlue3')
canvas.pack(expand = YES, fill = BOTH)

vidTB = Text(canvas,width=52,height=1)
vidTB.insert(INSERT, '' )
vidTB["bg"]="azure"
vidTB.bind("<Key>", lambda e: "break")
canvas.create_window(40,270,window=vidTB,anchor=NW)

lecTB = Text(canvas,width=52,height=1)
lecTB.insert(INSERT,'')
lecTB["bg"]="azure"
lecTB.bind("<Key>", lambda e: "break")
canvas.create_window(490,270,window=lecTB,anchor=NW)

tBimgW = Text(canvas,width=5,height=1)
tBimgW.insert(INSERT,str(imgW))
canvas.create_window(225,415,window=tBimgW ,anchor=NW)

tBnoC = Text(canvas,width=5,height=1)
tBnoC.insert(INSERT,str(noC))
canvas.create_window(225,465,window=tBnoC ,anchor=NW)

tBminPt = Text(canvas,width=5,height=1)
tBminPt.insert(INSERT,str(minPts))
canvas.create_window(225,515,window=tBminPt ,anchor=NW)

tBThres = Text(canvas,width=5,height=1)
tBThres.insert(INSERT,str(matchThres))
canvas.create_window(225,565,window=tBThres ,anchor=NW)

tBEstab = Text(canvas,width=5,height=1)
tBEstab.insert(INSERT,str(estabThres))
canvas.create_window(625,415,window=tBEstab ,anchor=NW)

tBDec= Text(canvas,width=5,height=1)
tBDec.insert(INSERT,str(decayThres))
canvas.create_window(625,465,window=tBDec ,anchor=NW)

tBFilter= Text(canvas,width=5,height=1)
tBFilter.insert(INSERT,str(filterMse))
canvas.create_window(625,515,window=tBFilter ,anchor=NW)

tBPPT= Text(canvas,width=30,height=1)
tBPPT.insert(INSERT,str(pptName))
canvas.create_window(335,510,window=tBPPT ,anchor=NW)

from PIL import Image
load = Image.open("tempImages/slidethru4.png")
render = ImageTk.PhotoImage(load)
img = Label(image=render)
img.image = render
img.place(x=0, y=0)

tb1=""

def canvasDisplay():
	global canvas,txtInS,root,tb1,tb1b,tb2,tbl1,tbl2,tbl3,tbl4,tbr1,tbr2,tbr3,tbr4,vidTB,tBnoC,\
	tBimgW,tBminPt,tBThres,tBEstab,tBDec,noC,imgW,matchThres, estabThres,decayThres, minPts,\
	filterMse,pptName, tbsl,tbr3
	#canvas.delete("all")

	if tb1 != "":
		canvas.delete(tb1)
		canvas.delete(tb2)
		canvas.delete(tb1b)
		canvas.delete(tbl1)
		canvas.delete(tbl2)
		canvas.delete(tbl3)
		canvas.delete(tbl4)
		canvas.delete(tbr1)
		canvas.delete(tbr2)
		canvas.delete(tbr3)
		canvas.delete(tbsl)

	canvas.create_line(0,400,900,400,fill="DeepSkyBlue4",width = 2)

	canvas.create_line(450,130,450,400,fill="DeepSkyBlue4",width = 2)
	canvas.create_line(0,130,900,130,fill="DeepSkyBlue4",width = 2)

	tb1=canvas.create_text(220,250,text="Video Path",fill="RoyalBlue4")
	tb1b=canvas.create_text(670,250,text="Folder Path",fill="RoyalBlue4")

	tbl1=canvas.create_text(165,425,text="Resize Width",fill="RoyalBlue4")
	tbl2=canvas.create_text(165,475,text="No of Cells",fill="RoyalBlue4")
	tbl3=canvas.create_text(165,525,text="Min Cell Density",fill="RoyalBlue4")
	tbl4=canvas.create_text(165,575,text="Match Thresh",fill="RoyalBlue4")

	tbr1=canvas.create_text(735,425,text="Establish Time",fill="RoyalBlue4")
	tbr2=canvas.create_text(735,475,text="Decay Time",fill="RoyalBlue4")
	tbr3=canvas.create_text(735,525,text="Filter Thresh",fill="RoyalBlue4")
	tbsl=canvas.create_text(445,490,text="PPT Name",fill="RoyalBlue4")

	displayButton("Generate Slides from Lecture Video",100,150,350,200,"White","DeepSkyBlue4")
	displayButton("Upload Lecture Video",100,325,350,375,"White","DeepSkyBlue4")
	displayButton("Generate Slides from Lesson Folder",550,150,800,200,"White","DeepSkyBlue4")
	displayButton("Select Lesson Folder",550,325,800,375,"White","DeepSkyBlue4")

	canvas.create_oval(450-15, 450-15, 450+15, 450+15, width=1, fill = "white")
	if txtInS:
		canvas.create_oval(450-10, 450-10, 450+10, 450+10, width=1, fill = "aquamarine4")
	else:
		canvas.create_oval(450-15, 450-15, 450+15, 450+15, width=1, fill = "white")

	tb2=canvas.create_text(450,420,text="Use Transcript",fill="RoyalBlue4")

	try:
		imgW = int(tBimgW.get(1.0,END))
		noC = int(tBnoC.get(1.0,END))
		minPts = int(tBminPt.get(1.0,END))
		matchThres = int(tBThres.get(1.0,END))
		estabThres = int(tBEstab.get(1.0,END))
		decayThres = int(tBDec.get(1.0,END))
		pptName = str(tBPPT.get(1.0,END))
		pptName = pptName[ : len(pptName)-1 ]
		filterMse = int(tBFilter.get(1.0,END))
	except:
		pass

def displayButton(txt,x1,y1,x2,y2,color,txtCol):
	global canvas
	canvas.create_rectangle(x1,y1,x2,y2, fill=color)
	canvas.create_text((x1+x2)/2,(y1+y2)/2,text=txt,fill=txtCol)

def mousePressed(event):
	global canvas,root,txtInS,argList,argList2,vidTB,lecTB,vloaded1,vloaded2
	eX = event.x
	eY = event.y
	#print(eX)
	#print(eY)
	a = (eX,eY)
	b = (450,450)
	dstRad = math.sqrt(sum( (a - b)**2 for a, b in zip(a, b)))
	if dstRad < 15:
		txtInS = not txtInS

	if eX>100 and eY>325 and eX<350 and eY<375:

		root.filename  = filedialog.askopenfilename(initialdir = "/",title = "Select file",filetypes = (("mp4 files","*.mp4"),("all files","*.*")))
		#print(root.filename)
		if 'Data/' in root.filename:
			filaname = root.filename[ root.filename.find('Data/')  : ]
			lstsl = filaname.rfind("/")
			folderP = filaname[:lstsl]
			fileP = filaname[lstsl+1:]
			argList = ['r',folderP,fileP]
			vidTB.delete('1.0', END)
			vidTB.insert(INSERT, argList[1]+"/"+argList[2]     )
			vloaded1=True

	if eX>550 and eY>325 and eX<800 and eY<375:

		folder = askdirectory()

		if 'Data/' in folder:
			folderP = folder[ folder.find('Data/') : ]
			argList2 = ['n',folderP]
			lecTB.delete('1.0', END)
			lecTB.insert(INSERT, folderP )
			vloaded2=True

	if eX>100 and eY>151 and eX<350 and eY<200:
		if vloaded1:
			videoConversion(1)

	if eX>552 and eY>151 and eX<800 and eY<200:
		if vloaded2:
			videoConversion(2)

def videoConversion(opt):
	l2d = Lec2Doc()
	l2d.runProcess(opt)

def updateCanvas():
	global canvas,root
	canvasDisplay()
	root.after(1,updateCanvas)

#####################################
canvas.bind("<Button-1>",mousePressed)
canvasDisplay()
updateCanvas()
root.mainloop()